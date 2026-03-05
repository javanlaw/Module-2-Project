from pptx import Presentation
from pptx.util import Inches, Pt

def create_netflix_presentation():
    prs = Presentation()

    # --- SLIDE 1: TITLE ---
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Modernizing Content Analytics: Netflix ELT Pipeline"
    slide.placeholders[1].text = "Building a Scalable Star Schema with dbt and BigQuery\nPresented by: Data Engineering Team"

    # --- SLIDE 2: THE BUSINESS CHALLENGE ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "The Business Challenge"
    content = slide.placeholders[1].text_frame
    content.text = "Current State: Manual handling of netflix_titles.csv"
    p = content.add_paragraph()
    p.text = "• Data Integrity: Duplicates in show_id skewing metrics."
    p = content.add_paragraph()
    p.text = "• Latency: High manual effort for cleaning and formatting."

    # --- SLIDE 3: TECHNICAL ARCHITECTURE ---
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Technical Architecture"
    content = slide.placeholders[1].text_frame
    content.text = "ELT Framework (Extract, Load, Transform)"
    for bullet in ["Ingestion: Python/Pandas automation to BigQuery.", 
                   "Compute: Google BigQuery as the SQL engine.",
                   "Transformation: dbt for modular SQL modeling."]:
        p = content.add_paragraph()
        p.text = f"• {bullet}"

    # --- SLIDE 4: DATA MODELING STRATEGY ---
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Star Schema Evolution"
    content = slide.placeholders[1].text_frame
    content.text = "Moving from Flat Files to Relational Models"
    for bullet in ["Staging Layer: stg_netflix handles deduplication.", 
                   "Fact Table: 'movies' (Quantitative data).",
                   "Dimensions: 'dim_directors' and 'dim_locations'."]:
        p = content.add_paragraph()
        p.text = f"• {bullet}"

    # --- SLIDE 5: BUSINESS VALUE ---
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Business Value Proposition"
    content = slide.placeholders[1].text_frame
    for bullet in ["Operational Efficiency: Automated ingestion reduces manual labor.", 
                   "Cost Optimization: Strategic materialization (Views vs. Tables).",
                   "Strategic Insights: Mapping regional growth and talent trends."]:
        p = content.add_paragraph()
        p.text = f"• {bullet}"

    # Save the presentation
    prs.save('Netflix_ELT_Project.pptx')
    print("Presentation created successfully: Netflix_ELT_Project.pptx")

if __name__ == "__main__":
    create_netflix_presentation()